"""Document Comparator v5 — Dual engine: AI Functions (SQL) or Local Markdown + Model Serving."""

import io
import json
import logging
import os
import time
import urllib.request
from pathlib import Path
from uuid import uuid4

from fastapi import FastAPI, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

app = FastAPI(title="Document Comparator", version="5.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

AVAILABLE_MODELS = [
    {"id": "databricks-meta-llama-3-3-70b-instruct", "name": "Llama 3.3 70B", "provider": "Meta"},
    {"id": "databricks-claude-sonnet-4-6", "name": "Claude Sonnet 4.6", "provider": "Anthropic"},
    {"id": "databricks-claude-opus-4-6", "name": "Claude Opus 4.6", "provider": "Anthropic"},
]

ENGINES = [
    {"id": "ai_functions", "name": "AI Functions (SQL)", "description": "ai_parse_document() + ai_query() via SQL Warehouse"},
    {"id": "direct_llm", "name": "Direct LLM", "description": "Local markdown conversion + Model Serving API"},
]

CATALOG = os.environ.get("CATALOG_NAME", "")
SCHEMA = os.environ.get("SCHEMA_NAME", "")
VOLUME = os.environ.get("VOLUME_NAME", "")
VOLUME_PATH = f"/Volumes/{CATALOG}/{SCHEMA}/{VOLUME}"
WAREHOUSE_ID = os.environ.get("DATABRICKS_WAREHOUSE_ID", "")

_ws_client = None

def get_ws():
    global _ws_client
    if _ws_client is None:
        from databricks.sdk import WorkspaceClient
        _ws_client = WorkspaceClient()
    return _ws_client


COMPARISON_PROMPT_TEMPLATE = """You are a document alignment analyst. Compare the two documents below.

IMPORTANT: Your response MUST use proper markdown formatting with ## headings, bullet lists with -, and **bold** text.

=== POWERPOINT ===
{pptx_text}

=== DOCUMENT ===
{doc_text}

Respond in this EXACT markdown format:

## Overall Similarity Score

**Score: X%** — [1-2 sentence explanation]

## Aligned Content

- **[Topic 1]** — [PowerPoint slide X] matches [Document section Y]. [Detail]
- **[Topic 2]** — [Detail with specific numbers]

## Divergences Found

- **[Divergence 1]**
  - PowerPoint says: [exact claim]
  - Document says: [exact claim]
  - Severity: **Major/Moderate/Minor**
  - Impact: [explanation]

## Content in PowerPoint Only

- [Item not found in the document]

## Content in Document Only

- [Item not found in the PowerPoint]

## Recommendation

- [Specific action 1]
- [Specific action 2]

Be precise with numbers, percentages, and exact claims. Do not generalize."""


# ══════════════════════════════════════════════════════════════
# ENGINE 1: AI Functions (SQL Warehouse)
# ══════════════════════════════════════════════════════════════

def run_sql(statement: str) -> list[dict]:
    """Execute SQL on the warehouse with async polling."""
    w = get_ws()
    from databricks.sdk.service.sql import StatementState

    resp = w.statement_execution.execute_statement(
        warehouse_id=WAREHOUSE_ID, statement=statement, wait_timeout="0s"
    )
    statement_id = resp.statement_id
    max_wait, elapsed, poll_interval = 600, 0, 2

    while elapsed < max_wait:
        status = w.statement_execution.get_statement(statement_id)
        if status.status.state == StatementState.SUCCEEDED:
            if not status.result or not status.result.data_array:
                return []
            columns = [col.name for col in status.manifest.schema.columns]
            return [dict(zip(columns, row)) for row in status.result.data_array]
        if status.status.state in (StatementState.FAILED, StatementState.CANCELED, StatementState.CLOSED):
            error_msg = status.status.error.message if status.status.error else "Unknown SQL error"
            raise RuntimeError(f"SQL failed: {error_msg}")
        time.sleep(poll_interval)
        elapsed += poll_interval
        if elapsed > 30:
            poll_interval = 5

    raise RuntimeError(f"SQL query timed out after {max_wait}s")


async def upload_to_volume(filename: str, content: bytes) -> str:
    w = get_ws()
    safe_name = f"{uuid4().hex[:8]}_{filename.replace(' ', '_')}"
    file_path = f"{VOLUME_PATH}/{safe_name}"
    try:
        w.files.upload(file_path, io.BytesIO(content), overwrite=True)
        return file_path
    except Exception as e:
        if "NOT_FOUND" in str(e) or "RESOURCE_DOES_NOT_EXIST" in str(e):
            try:
                run_sql(f"CREATE VOLUME IF NOT EXISTS {CATALOG}.{SCHEMA}.{VOLUME}")
            except:
                pass
            w.files.upload(file_path, io.BytesIO(content), overwrite=True)
            return file_path
        raise


def cleanup_volume_file(path: str):
    try:
        get_ws().files.delete(path)
    except:
        pass


async def compare_ai_functions(pptx_path: str, doc_path: str, model_id: str, doc_name: str = "") -> dict:
    """ai_parse_document() + ai_query() via SQL. Uses openpyxl for xlsx since ai_parse_document doesn't support it."""
    ep = pptx_path.replace("'", "''")
    is_xlsx = doc_name.endswith(".xlsx")

    prompt_template = (
        "'You are a document alignment analyst. Compare the two documents below.\\n\\n"
        "IMPORTANT: Use proper markdown with ## headings, - bullets, **bold**.\\n\\n"
        "=== POWERPOINT ===\\n', pptx.content, '\\n\\n"
        "=== DOCUMENT ===\\n', {doc_col}, '\\n\\n"
        "Respond in this EXACT format:\\n\\n"
        "## Overall Similarity Score\\n\\n**Score: X%%** — [explanation]\\n\\n"
        "## Aligned Content\\n\\n- **[Topic]** — [slide X] matches [section Y]\\n\\n"
        "## Divergences Found\\n\\n- **[Divergence]**\\n  - PowerPoint says: [claim]\\n  - Document says: [claim]\\n  - Severity: **Major/Moderate/Minor**\\n\\n"
        "## Content in PowerPoint Only\\n\\n- [items]\\n\\n"
        "## Content in Document Only\\n\\n- [items]\\n\\n"
        "## Recommendation\\n\\n- [actions]\\n\\nBe precise with numbers.'"
    )

    if is_xlsx:
        doc_md = xlsx_to_markdown(get_ws().files.download(doc_path).contents.read())
        esc_doc = doc_md.replace("'", "''").replace("\\", "\\\\")
        doc_col_ref = f"'{esc_doc}'"
        prompt_sql = prompt_template.replace("{doc_col}", doc_col_ref)
        sql = f"""
        WITH pptx_parsed AS (
          SELECT ai_parse_document(content, map('version', '2.0', 'descriptionElementTypes', '*'))::STRING AS content
          FROM read_files('{ep}', format => 'binaryFile')
        )
        SELECT pptx.content AS pptx_text, {doc_col_ref} AS doc_text,
          ai_query('{model_id}', CONCAT({prompt_sql})) AS report
        FROM pptx_parsed pptx
        """
    else:
        ed = doc_path.replace("'", "''")
        prompt_sql = prompt_template.replace("{doc_col}", "doc.content")
        sql = f"""
        WITH pptx_parsed AS (
          SELECT ai_parse_document(content, map('version', '2.0', 'descriptionElementTypes', '*'))::STRING AS content
          FROM read_files('{ep}', format => 'binaryFile')
        ),
        doc_parsed AS (
          SELECT ai_parse_document(content, map('version', '2.0', 'descriptionElementTypes', '*'))::STRING AS content
          FROM read_files('{ed}', format => 'binaryFile')
        )
        SELECT pptx.content AS pptx_text, doc.content AS doc_text,
          ai_query('{model_id}', CONCAT({prompt_sql})) AS report
        FROM pptx_parsed pptx CROSS JOIN doc_parsed doc
        """

    rows = run_sql(sql)
    if not rows:
        raise RuntimeError("Query returned no results")
    return {"pptx_text": rows[0].get("pptx_text", ""), "doc_text": rows[0].get("doc_text", ""), "report": rows[0].get("report", "")}


# ══════════════════════════════════════════════════════════════
# ENGINE 2: Local Markdown + Direct Model Serving
# ══════════════════════════════════════════════════════════════

def pptx_to_markdown(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    md = ["# PowerPoint Presentation\n"]
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                    title = shape.text_frame.text.strip()
                else:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            body.append(t)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                if rows:
                    md.append("| " + " | ".join(rows[0]) + " |")
                    md.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                    for r in rows[1:]:
                        md.append("| " + " | ".join(r) + " |")
        md.append(f"## Slide {i}: {title or f'Slide {i}'}\n")
        for b in body:
            md.append(f"- {b}")
        md.append("")
    return "\n".join(md)


def docx_to_markdown(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    md = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        s = para.style.name
        if s == "Title": md.append(f"# {t}\n")
        elif s == "Heading 1": md.append(f"## {t}\n")
        elif s == "Heading 2": md.append(f"### {t}\n")
        elif "List" in s or "Bullet" in s: md.append(f"- {t}")
        else: md.append(f"{t}\n")
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if rows:
            md.append("| " + " | ".join(rows[0]) + " |")
            md.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
            for r in rows[1:]:
                md.append("| " + " | ".join(r) + " |")
    return "\n".join(md)


def xlsx_to_markdown(content: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), data_only=True)
    md = ["# Excel Spreadsheet\n"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md.append(f"## Sheet: {sheet_name}\n")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            md.append("(empty sheet)\n")
            continue
        # Header
        headers = [str(c) if c is not None else "" for c in rows[0]]
        md.append("| " + " | ".join(headers) + " |")
        md.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in rows[1:]:
            vals = [str(c) if c is not None else "" for c in row]
            if any(v for v in vals):
                md.append("| " + " | ".join(vals) + " |")
        md.append("")
    return "\n".join(md)


def txt_to_markdown(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def delta_table_to_markdown(table_name: str) -> str:
    """Query a Delta table via SQL and format as markdown."""
    rows = run_sql(f"SELECT * FROM {table_name} LIMIT 500")
    if not rows:
        return f"# Delta Table: {table_name}\n\n(empty table)"
    columns = list(rows[0].keys())
    md = [f"# Delta Table: {table_name}\n"]
    md.append("| " + " | ".join(columns) + " |")
    md.append("| " + " | ".join(["---"] * len(columns)) + " |")
    for row in rows:
        vals = [str(row.get(c, "")) if row.get(c) is not None else "" for c in columns]
        md.append("| " + " | ".join(vals) + " |")
    return "\n".join(md)


def parse_to_markdown(filename: str, content: bytes) -> str:
    if filename.endswith(".pptx"): return pptx_to_markdown(content)
    elif filename.endswith(".docx"): return docx_to_markdown(content)
    elif filename.endswith(".xlsx"): return xlsx_to_markdown(content)
    elif filename.endswith(".txt"): return txt_to_markdown(content)
    else: raise ValueError(f"Unsupported: {filename}")


async def call_model_serving(pptx_md: str, doc_md: str, model_id: str) -> str:
    """Call Model Serving REST API directly."""
    w = get_ws()
    host = w.config.host.rstrip("/")
    headers = w.config.authenticate()
    auth_header = headers.get("Authorization", "")

    prompt = COMPARISON_PROMPT_TEMPLATE.format(pptx_text=pptx_md, doc_text=doc_md)

    payload = json.dumps({
        "messages": [
            {"role": "system", "content": "You are an expert document comparison analyst. Always respond in well-structured markdown."},
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 4000,
        "temperature": 0.1
    }).encode()

    url = f"{host}/serving-endpoints/{model_id}/invocations"
    req = urllib.request.Request(url, data=payload,
        headers={"Authorization": auth_header, "Content-Type": "application/json"})

    with urllib.request.urlopen(req, timeout=300) as resp:
        result = json.loads(resp.read())
        return result["choices"][0]["message"]["content"]


async def compare_direct_llm(pptx_content: bytes, pptx_name: str, doc_content: bytes, doc_name: str, model_id: str) -> dict:
    """Parse locally to markdown, then call Model Serving directly."""
    pptx_md = parse_to_markdown(pptx_name, pptx_content)
    doc_md = parse_to_markdown(doc_name, doc_content)
    report = await call_model_serving(pptx_md, doc_md, model_id)
    return {"pptx_text": pptx_md, "doc_text": doc_md, "report": report}


# ══════════════════════════════════════════════════════════════
# API Routes
# ══════════════════════════════════════════════════════════════

@app.get("/api/models")
def list_models():
    return {"models": AVAILABLE_MODELS}


@app.get("/api/engines")
def list_engines():
    return {"engines": ENGINES}


@app.post("/api/compare")
async def compare_documents(
    pptx_file: UploadFile = File(...),
    doc_file: UploadFile = File(None),
    model: str = Form("databricks-meta-llama-3-3-70b-instruct"),
    engine: str = Form("direct_llm"),
    delta_table: str = Form("")
):
    """Compare documents using the selected engine and model. Supports docx, xlsx, pdf, txt, or Delta table."""
    pptx_content = await pptx_file.read()

    pptx_name = pptx_file.filename or "presentation.pptx"
    if not pptx_name.endswith(".pptx"):
        return JSONResponse(status_code=400, content={"error": "First file must be a .pptx PowerPoint file"})

    # Determine comparison target
    use_delta = bool(delta_table.strip())
    if not use_delta:
        if doc_file is None:
            return JSONResponse(status_code=400, content={"error": "Provide either a document file or a Delta table name"})
        doc_content = await doc_file.read()
        doc_name = doc_file.filename or "document"
        if not doc_name.endswith((".docx", ".xlsx", ".pdf", ".txt")):
            return JSONResponse(status_code=400, content={"error": "Second file must be .docx, .xlsx, .pdf, or .txt"})
    else:
        doc_content = None
        doc_name = delta_table.strip()

    pptx_path = None
    doc_path = None

    try:
        if use_delta:
            # Delta table comparison — query the table, format as markdown, compare
            logger.info(f"Engine: {engine} | Model: {model} | Delta: {doc_name}")
            delta_md = delta_table_to_markdown(doc_name)

            if engine == "ai_functions":
                pptx_path = await upload_to_volume(pptx_name, pptx_content)
                pptx_parsed = run_sql(f"""
                    SELECT ai_parse_document(content, map('version', '2.0', 'descriptionElementTypes', '*'))::STRING AS parsed
                    FROM read_files('{pptx_path.replace("'", "''")}', format => 'binaryFile')
                """)
                pptx_text = pptx_parsed[0]["parsed"] if pptx_parsed else ""
                escaped_prompt = COMPARISON_PROMPT_TEMPLATE.format(pptx_text=pptx_text, doc_text=delta_md).replace("'", "''")
                report_rows = run_sql(f"SELECT ai_query('{model}', '{escaped_prompt}') AS report")
                result = {"pptx_text": pptx_text, "doc_text": delta_md, "report": report_rows[0]["report"] if report_rows else ""}
            else:
                pptx_md = parse_to_markdown(pptx_name, pptx_content)
                report = await call_model_serving(pptx_md, delta_md, model)
                result = {"pptx_text": pptx_md, "doc_text": delta_md, "report": report}

        elif engine == "ai_functions":
            logger.info(f"Engine: AI Functions | Model: {model}")
            pptx_path = await upload_to_volume(pptx_name, pptx_content)
            doc_path = await upload_to_volume(doc_name, doc_content)
            result = await compare_ai_functions(pptx_path, doc_path, model, doc_name)

        elif engine == "direct_llm":
            logger.info(f"Engine: Direct LLM | Model: {model}")
            result = await compare_direct_llm(pptx_content, pptx_name, doc_content, doc_name, model)

        else:
            return JSONResponse(status_code=400, content={"error": f"Unknown engine: {engine}"})

        logger.info(f"Done: pptx={len(result['pptx_text'])} chars, doc={len(result['doc_text'])} chars, report={len(result['report'])} chars")

        def extract_sections(text, max_sections=10):
            sections = []
            current_title = "Content"
            current_content = []
            for line in text.split("\n"):
                s = line.strip()
                if s and (s.startswith("#") or (s.isupper() and len(s) < 80)):
                    if current_content:
                        sections.append({"title": current_title, "preview": " ".join(current_content)[:200]})
                    current_title = s.lstrip("#").strip()
                    current_content = []
                elif s:
                    current_content.append(s)
            if current_content:
                sections.append({"title": current_title, "preview": " ".join(current_content)[:200]})
            return sections[:max_sections]

        pptx_text = result["pptx_text"]
        doc_text = result["doc_text"]

        return {
            "pptx_file": pptx_name, "doc_file": doc_name,
            "model": model, "engine": engine,
            "pptx_sections": len(extract_sections(pptx_text)),
            "doc_sections": len(extract_sections(doc_text)),
            "pptx_content": extract_sections(pptx_text),
            "doc_content": extract_sections(doc_text),
            "pptx_markdown": pptx_text, "doc_markdown": doc_text,
            "comparison_report": result["report"]
        }

    except Exception as e:
        logger.error(f"Comparison failed: {e}", exc_info=True)
        return JSONResponse(status_code=500, content={"error": str(e)})
    finally:
        if pptx_path: cleanup_volume_file(pptx_path)
        if doc_path: cleanup_volume_file(doc_path)


@app.get("/api/health")
def health():
    return {"status": "ok", "version": "5.0", "engines": ["ai_functions", "direct_llm"]}


# ── Static Files & SPA ──────────────────────────────────────

dist = Path(__file__).parent / "frontend" / "dist"
if dist.exists():
    app.mount("/assets", StaticFiles(directory=dist / "assets"), name="assets")

    @app.get("/{full_path:path}")
    async def serve_spa(full_path: str):
        return FileResponse(dist / "index.html")
else:
    @app.get("/")
    def root():
        return {"message": "Document Comparator API v5.0. POST /api/compare with engine=ai_functions|direct_llm"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
