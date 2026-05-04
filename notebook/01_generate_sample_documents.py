# Databricks notebook source
# MAGIC %md
# MAGIC # 01 — Generate Sample Documents
# MAGIC
# MAGIC Generates sample documents for the comparison demo — PowerPoint, Word, Excel, and a Delta table.
# MAGIC All cover the same Q1 2026 earnings material with deliberate differences.
# MAGIC
# MAGIC **Output:**
# MAGIC - `quarterly_review.pptx` — Earnings presentation (6 slides)
# MAGIC - `earnings_summary.docx` — Earnings write-up (7 sections)
# MAGIC - `earnings_data.xlsx` — Earnings financials spreadsheet (4 sheets)
# MAGIC - Delta table: `earnings_summary` — Structured earnings data

# COMMAND ----------

# MAGIC %pip install python-pptx python-docx openpyxl --quiet
# MAGIC dbutils.library.restartPython()

# COMMAND ----------

dbutils.widgets.text("catalog", "ronguerrero", "Catalog")
dbutils.widgets.text("schema", "capital_markets_ai", "Schema")

CATALOG = dbutils.widgets.get("catalog")
SCHEMA = dbutils.widgets.get("schema")
VOLUME_PATH = f"/Volumes/{CATALOG}/{SCHEMA}/document_comparison"
TABLE = f"{CATALOG}.{SCHEMA}.earnings_summary"

spark.sql(f"CREATE SCHEMA IF NOT EXISTS {CATALOG}.{SCHEMA}")
spark.sql(f"CREATE VOLUME IF NOT EXISTS {CATALOG}.{SCHEMA}.document_comparison")
spark.sql(f"USE CATALOG {CATALOG}")

print(f"Catalog: {CATALOG}")
print(f"Schema:  {SCHEMA}")
print(f"Volume:  {VOLUME_PATH}")
print(f"Table:   {TABLE}")

# COMMAND ----------

# MAGIC %md
# MAGIC ### Generate PowerPoint Presentation

# COMMAND ----------

import shutil, tempfile
from pptx import Presentation

prs = Presentation()
slides = [
    ("Q1 2026 Business Review", [
        "Revenue grew 23% YoY to $4.82B",
        "Cloud Platform segment leads with 31% growth",
        "Operating margin expanded to 30.1%",
        "Raised FY2026 guidance to $20.5B-$21.0B"
    ]),
    ("Revenue Breakdown", [
        "Cloud Platform: $2.41B (50% of total, +31% YoY)",
        "AI & Analytics: $1.21B (25% of total, +29% YoY)",
        "Cybersecurity: $723M (15% of total, +15% YoY)",
        "Professional Services: $482M (10% of total, +8% YoY)"
    ]),
    ("Key Metrics", [
        "Gross Margin: 70.1% (up 210bps YoY)",
        "Free Cash Flow: $1.62B (+41% YoY)",
        "Net Income: $1.18B (+45% YoY)",
        "EPS: $2.94 (beat consensus by 9.7%)"
    ]),
    ("Strategic Priorities", [
        "Accelerate AI platform adoption across enterprise customers",
        "Expand international presence, particularly in APAC and EMEA",
        "Drive cloud migration for existing on-premise customers",
        "Invest in cybersecurity capabilities through M&A"
    ]),
    ("Risk Factors", [
        "Macroeconomic slowdown could impact enterprise IT spending",
        "Competitive pressure from hyperscalers (AWS, Azure, GCP)",
        "Foreign currency headwinds from strong USD (~150bps impact)",
        "Talent retention costs in AI/ML engineering"
    ]),
    ("FY2026 Outlook", [
        "Q2 Revenue: $5.05B - $5.15B (21-23% YoY growth)",
        "FY2026 Revenue: $20.5B - $21.0B (raised from $19.8B-$20.3B)",
        "Operating Margin target: 31-32%",
        "FY2026 EPS: $12.80 - $13.20"
    ]),
]
for title, bullets in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame; tf.clear()
    for i, b in enumerate(bullets):
        if i == 0: tf.paragraphs[0].text = b
        else: tf.add_paragraph().text = b

tmp_pptx = tempfile.mktemp(suffix=".pptx")
prs.save(tmp_pptx)
shutil.copy(tmp_pptx, f"{VOLUME_PATH}/quarterly_review.pptx")
print(f"Created: quarterly_review.pptx ({len(slides)} slides)")

# COMMAND ----------

# MAGIC %md
# MAGIC ### Generate Word Document

# COMMAND ----------

import shutil, tempfile
from docx import Document

doc = Document()
doc.add_heading('TechNova Corporation — Q1 2026 Earnings Summary', 0)

doc.add_heading('Financial Highlights', level=1)
doc.add_paragraph('TechNova reported strong Q1 2026 results with total revenue of $4.82 billion, representing 23.3% year-over-year growth.')
for item in ['Revenue: $4.82B vs. $4.58B consensus (+5.2% beat)', 'Gross Margin: 70.1%, expanding 210bps', 'Operating Income: $1.45B with 30.1% margin', 'Net Income: $1.18B, up 45.3% YoY', 'Diluted EPS: $2.94 vs. $2.68 consensus (+9.7% beat)', 'Free Cash Flow: $1.62B, 33.6% FCF margin']:
    doc.add_paragraph(item, style='List Bullet')

doc.add_heading('Segment Performance', level=1)
doc.add_paragraph('Cloud Platform: $2.41B (50%, +31.2% YoY). AI & Analytics: $1.21B (25%, +28.7%). Cybersecurity: $723M (15%, +15.4%). Professional Services: $482M (10%, +8.1%).')

doc.add_heading('Strategic Initiatives', level=1)
for item in ['AI Platform Expansion: New GenAI capabilities', 'Geographic Growth: APAC +40% headcount, EMEA +25%', 'Cloud Migration: Converting 200+ on-premise customers', 'Cybersecurity M&A: Evaluating 3 acquisition targets', 'Partner Ecosystem: 15 new ISV integrations by Q4']:
    doc.add_paragraph(item, style='List Bullet')

doc.add_heading('Risk Assessment', level=1)
doc.add_paragraph('Key risks: IT spending slowdown, hyperscaler competition, 150bps FX headwind, customer concentration (top 10 = 28% revenue), EU AI Act regulatory risk.')

doc.add_heading('Forward Guidance', level=1)
for item in ['Q2 Revenue: $5.05B-$5.15B (21-23% YoY)', 'FY2026 Revenue: $20.5B-$21.0B (raised)', 'Operating Margin: 31-32%', 'EPS: $12.80-$13.20', 'Capital Allocation: $2.5B share repurchase']:
    doc.add_paragraph(item, style='List Bullet')

doc.add_heading('Analyst Consensus', level=1)
doc.add_paragraph('8 analysts raised price targets, avg $142 (+18% upside). Consensus: Strong Buy (12/14 analysts).')

tmp_docx = tempfile.mktemp(suffix=".docx")
doc.save(tmp_docx)
shutil.copy(tmp_docx, f"{VOLUME_PATH}/earnings_summary.docx")
print("Created: earnings_summary.docx")

# COMMAND ----------

# MAGIC %md
# MAGIC ### Generate Excel Spreadsheet

# COMMAND ----------

import shutil, tempfile
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

wb = Workbook()

header_font = Font(bold=True, color="FFFFFF", size=11)
header_fill = PatternFill(start_color="0A1628", end_color="0A1628", fill_type="solid")
border = Border(bottom=Side(style='thin', color='CCCCCC'))

# Sheet 1: Financial Summary
ws = wb.active; ws.title = "Financial Summary"
for col, h in enumerate(["Metric", "Q1 2026", "Q1 2025", "YoY Change"], 1):
    cell = ws.cell(row=1, column=col, value=h); cell.font = header_font; cell.fill = header_fill
for r, row_data in enumerate([
    ["Revenue", "$4.82B", "$3.91B", "+23.3%"], ["Gross Profit", "$3.38B", "$2.66B", "+27.1%"],
    ["Operating Income", "$1.45B", "$1.02B", "+42.2%"], ["Net Income", "$1.18B", "$812M", "+45.3%"],
    ["Diluted EPS", "$2.94", "$2.01", "+46.3%"], ["Free Cash Flow", "$1.62B", "$1.15B", "+40.9%"],
    ["Gross Margin", "70.1%", "68.0%", "+210 bps"], ["Operating Margin", "30.1%", "26.1%", "+400 bps"],
], 2):
    for c, val in enumerate(row_data, 1): ws.cell(row=r, column=c, value=val).border = border
for col in range(1, 5): ws.column_dimensions[chr(64+col)].width = 20

# Sheet 2: Revenue by Segment
ws2 = wb.create_sheet("Revenue by Segment")
for col, h in enumerate(["Segment", "Revenue", "% of Total", "YoY Growth"], 1):
    cell = ws2.cell(row=1, column=col, value=h); cell.font = header_font; cell.fill = header_fill
for r, row_data in enumerate([["Cloud Platform", "$2.41B", "50%", "+31.2%"], ["AI & Analytics", "$1.21B", "25%", "+28.7%"],
    ["Cybersecurity", "$723M", "15%", "+15.4%"], ["Professional Services", "$482M", "10%", "+8.1%"]], 2):
    for c, val in enumerate(row_data, 1): ws2.cell(row=r, column=c, value=val)

# Sheet 3: Guidance
ws3 = wb.create_sheet("Guidance")
for col, h in enumerate(["Metric", "Low", "High", "Prior Low", "Prior High"], 1):
    cell = ws3.cell(row=1, column=col, value=h); cell.font = header_font; cell.fill = header_fill
for r, row_data in enumerate([["Q2 2026 Revenue", "$5.05B", "$5.15B", "", ""],
    ["FY2026 Revenue", "$20.5B", "$21.0B", "$19.8B", "$20.3B"],
    ["FY2026 Operating Margin", "31%", "32%", "30%", "31%"],
    ["FY2026 EPS", "$12.80", "$13.20", "$12.20", "$12.60"]], 2):
    for c, val in enumerate(row_data, 1): ws3.cell(row=r, column=c, value=val)

# Sheet 4: Risk Factors
ws4 = wb.create_sheet("Risk Factors")
for col, h in enumerate(["Risk Factor", "Severity", "Impact"], 1):
    cell = ws4.cell(row=1, column=col, value=h); cell.font = header_font; cell.fill = header_fill
for r, row_data in enumerate([
    ["Macroeconomic IT spending slowdown", "High", "Could reduce deal velocity 15-20%"],
    ["Hyperscaler competition", "High", "Pricing pressure on cloud segment"],
    ["Foreign currency headwinds", "Medium", "~150bps revenue impact"],
    ["Talent retention in AI/ML", "Medium", "Compensation costs rising 12% YoY"],
    ["Customer concentration", "Medium", "Top 10 clients = 28% of revenue"],
    ["EU AI Act compliance", "Low", "May delay features 1-2 quarters"]], 2):
    for c, val in enumerate(row_data, 1): ws4.cell(row=r, column=c, value=val)
ws4.column_dimensions['A'].width = 45; ws4.column_dimensions['C'].width = 50

tmp_xlsx = tempfile.mktemp(suffix=".xlsx")
wb.save(tmp_xlsx)
shutil.copy(tmp_xlsx, f"{VOLUME_PATH}/earnings_data.xlsx")
print("Created: earnings_data.xlsx (4 sheets)")

# COMMAND ----------

# MAGIC %md
# MAGIC ### Create Delta Table

# COMMAND ----------

spark.sql(f"""CREATE OR REPLACE TABLE {TABLE} (
  category STRING COMMENT 'Data category',
  metric STRING COMMENT 'Metric name',
  value STRING COMMENT 'Metric value',
  comparison STRING COMMENT 'Prior period or benchmark',
  change STRING COMMENT 'Change vs comparison',
  notes STRING COMMENT 'Additional context'
) COMMENT 'TechNova Q1 2026 earnings summary for document comparison demo'""")

spark.sql(f"""INSERT INTO {TABLE} VALUES
('Financial', 'Revenue', '$4.82B', '$3.91B (Q1 2025)', '+23.3% YoY', 'Beat consensus $4.58B by 5.2%'),
('Financial', 'Gross Profit', '$3.38B', '$2.66B (Q1 2025)', '+27.1% YoY', 'Gross margin 70.1%, up 210bps'),
('Financial', 'Operating Income', '$1.45B', '$1.02B (Q1 2025)', '+42.2% YoY', 'Operating margin 30.1%, up 400bps'),
('Financial', 'Net Income', '$1.18B', '$812M (Q1 2025)', '+45.3% YoY', NULL),
('Financial', 'Diluted EPS', '$2.94', '$2.01 (Q1 2025)', '+46.3% YoY', 'Beat consensus $2.68 by 9.7%'),
('Financial', 'Free Cash Flow', '$1.62B', '$1.15B (Q1 2025)', '+40.9% YoY', '33.6% FCF margin'),
('Segment', 'Cloud Platform', '$2.41B', '50% of total', '+31.2% YoY', 'Primary growth driver'),
('Segment', 'AI & Analytics', '$1.21B', '25% of total', '+28.7% YoY', NULL),
('Segment', 'Cybersecurity', '$723M', '15% of total', '+15.4% YoY', NULL),
('Segment', 'Professional Services', '$482M', '10% of total', '+8.1% YoY', NULL),
('Guidance', 'Q2 2026 Revenue', '$5.05B - $5.15B', NULL, '21-23% YoY growth', NULL),
('Guidance', 'FY2026 Revenue', '$20.5B - $21.0B', '$19.8B - $20.3B (prior)', 'Raised', NULL),
('Guidance', 'FY2026 Operating Margin', '31% - 32%', '30% - 31% (prior)', 'Raised', NULL),
('Guidance', 'FY2026 EPS', '$12.80 - $13.20', '$12.20 - $12.60 (prior)', 'Raised', NULL),
('Guidance', 'Share Repurchase', '$2.5B', NULL, 'New program', NULL),
('Strategy', 'AI Platform Adoption', 'Accelerating enterprise GenAI adoption', NULL, NULL, NULL),
('Strategy', 'International Expansion', 'APAC +40% headcount, EMEA +25%', NULL, NULL, NULL),
('Strategy', 'Cloud Migration', 'Converting 200+ on-premise customers', NULL, NULL, NULL),
('Strategy', 'Cybersecurity M&A', 'Evaluating 3 acquisition targets', NULL, NULL, NULL),
('Strategy', 'Partner Ecosystem', '15 new ISV integrations by Q4', NULL, NULL, 'Not in PowerPoint'),
('Risk', 'IT Spending Slowdown', 'Macroeconomic uncertainty', 'High', NULL, NULL),
('Risk', 'Hyperscaler Competition', 'AWS, Azure, GCP pricing pressure', 'High', NULL, NULL),
('Risk', 'FX Headwinds', '~150bps revenue impact from strong USD', 'Medium', NULL, NULL),
('Risk', 'Talent Retention', 'AI/ML engineering compensation rising 12% YoY', 'Medium', NULL, 'Not in PowerPoint'),
('Risk', 'Customer Concentration', 'Top 10 clients = 28% of revenue', 'Medium', NULL, 'Not in PowerPoint'),
('Risk', 'EU AI Act', 'May delay product features by 1-2 quarters', 'Low', NULL, 'Not in PowerPoint'),
('Analyst', 'Price Target Revisions', '8 analysts raised targets', 'Avg $142', '+18% upside', 'Not in PowerPoint'),
('Analyst', 'Rating', 'Strong Buy', '12 out of 14 analysts', NULL, 'Not in PowerPoint')
""")

print("Delta table created and populated")

# COMMAND ----------

display(spark.sql(f"SELECT category, COUNT(*) AS rows FROM {TABLE} GROUP BY category ORDER BY category"))

# COMMAND ----------

# MAGIC %md
# MAGIC ### Verify All Outputs

# COMMAND ----------

display(spark.sql(f"SELECT regexp_extract(path, '([^/]+)$', 1) AS filename, length(content) AS file_size_bytes FROM read_files('{VOLUME_PATH}/', format => 'binaryFile') ORDER BY filename"))

# COMMAND ----------

# MAGIC %md
# MAGIC **Created 4 comparison targets:**
# MAGIC
# MAGIC | Target | Type | What It Contains |
# MAGIC |--------|------|-----------------|
# MAGIC | `quarterly_review.pptx` | PowerPoint | 6 slides — the "source of truth" presentation |
# MAGIC | `earnings_summary.docx` | Word | 7 sections — narrative write-up with extra detail |
# MAGIC | `earnings_data.xlsx` | Excel | 4 sheets — financial tables, segments, guidance, risks |
# MAGIC | `earnings_summary` table | Delta | 30 rows — structured data with categories and notes |
# MAGIC
# MAGIC **Next:** Run `02_compare_documents` to compare the PowerPoint against any of these targets.
