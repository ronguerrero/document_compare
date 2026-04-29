# Databricks notebook source
# MAGIC %md
# MAGIC # Document Comparator: PowerPoint vs Document Alignment Check
# MAGIC
# MAGIC This notebook uses **`ai_parse_document()`** and **`ai_query()`** to parse and compare
# MAGIC a PowerPoint presentation against a document — no python-pptx or python-docx needed.
# MAGIC
# MAGIC **How it works:**
# MAGIC 1. Upload `.pptx` and `.docx` files to a Unity Catalog Volume
# MAGIC 2. `ai_parse_document()` extracts structured text from both files
# MAGIC 3. `ai_query()` compares content and generates an alignment report
# MAGIC
# MAGIC **Requirements:**
# MAGIC - Databricks Runtime 15.4+
# MAGIC - SQL Warehouse with Foundation Model APIs enabled
# MAGIC - Unity Catalog Volume for file storage
# MAGIC
# MAGIC **No library installs required** — everything runs as SQL functions.

# COMMAND ----------

# MAGIC %md
# MAGIC ## Configuration

# COMMAND ----------

CATALOG = "lakedoom_demo_catalog"
SCHEMA = "capital_markets_ai"
VOLUME = "document_comparison"
VOLUME_PATH = f"/Volumes/{CATALOG}/{SCHEMA}/{VOLUME}"

spark.sql(f"CREATE VOLUME IF NOT EXISTS {CATALOG}.{SCHEMA}.{VOLUME}")
print(f"Volume ready: {VOLUME_PATH}")
print(f"Upload your .pptx and .docx files to: {VOLUME_PATH}")

# COMMAND ----------

# MAGIC %md
# MAGIC ## Step 1: Generate Sample Documents (for demo)
# MAGIC
# MAGIC Skip this cell if you've uploaded your own files.

# COMMAND ----------

# MAGIC %pip install python-pptx python-docx --quiet
# MAGIC dbutils.library.restartPython()

# COMMAND ----------

CATALOG = "lakedoom_demo_catalog"
SCHEMA = "capital_markets_ai"
VOLUME = "document_comparison"
VOLUME_PATH = f"/Volumes/{CATALOG}/{SCHEMA}/{VOLUME}"

def create_sample_pptx(path):
    from pptx import Presentation
    prs = Presentation()
    slides = [
        ("Q1 2026 Business Review", [
            "Revenue grew 23% YoY to $4.82B",
            "Cloud Platform segment leads with 31% growth",
            "Operating margin expanded to 30.1%",
            "Raised FY2026 guidance to $20.5B-$21.0B"
        ]),
        ("Revenue Breakdown", [
            "Cloud Platform: $2.41B (50% of total, +31% YoY)",
            "AI & Analytics: $1.21B (25% of total, +29% YoY)",
            "Cybersecurity: $723M (15% of total, +15% YoY)",
            "Professional Services: $482M (10% of total, +8% YoY)"
        ]),
        ("Key Metrics", [
            "Gross Margin: 70.1% (up 210bps YoY)",
            "Free Cash Flow: $1.62B (+41% YoY)",
            "Net Income: $1.18B (+45% YoY)",
            "EPS: $2.94 (beat consensus by 9.7%)"
        ]),
        ("Strategic Priorities", [
            "Accelerate AI platform adoption across enterprise customers",
            "Expand international presence, particularly in APAC and EMEA",
            "Drive cloud migration for existing on-premise customers",
            "Invest in cybersecurity capabilities through M&A"
        ]),
        ("Risk Factors", [
            "Macroeconomic slowdown could impact enterprise IT spending",
            "Competitive pressure from hyperscalers (AWS, Azure, GCP)",
            "Foreign currency headwinds from strong USD (~150bps impact)",
            "Talent retention costs in AI/ML engineering"
        ]),
        ("FY2026 Outlook", [
            "Q2 Revenue: $5.05B - $5.15B (21-23% YoY growth)",
            "FY2026 Revenue: $20.5B - $21.0B (raised from $19.8B-$20.3B)",
            "Operating Margin target: 31-32%",
            "FY2026 EPS: $12.80 - $13.20"
        ]),
    ]
    for title, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame; tf.clear()
        for i, b in enumerate(bullets):
            if i == 0: tf.paragraphs[0].text = b
            else: tf.add_paragraph().text = b
    prs.save(path)
    print(f"Created: {path}")

def create_sample_docx(path):
    from docx import Document
    doc = Document()
    doc.add_heading('TechNova Corporation — Q1 2026 Earnings Summary', 0)

    doc.add_heading('Financial Highlights', level=1)
    doc.add_paragraph('TechNova reported strong Q1 2026 results with total revenue of $4.82 billion, representing 23.3% year-over-year growth. The company exceeded analyst expectations across all major metrics.')
    for item in ['Revenue: $4.82B vs. $4.58B consensus estimate (+5.2% beat)', 'Gross Margin: 70.1%, expanding 210 basis points from prior year', 'Operating Income: $1.45B with 30.1% operating margin', 'Net Income: $1.18B, up 45.3% year-over-year', 'Diluted EPS: $2.94 vs. $2.68 consensus (+9.7% beat)', 'Free Cash Flow: $1.62B, representing 33.6% FCF margin']:
        doc.add_paragraph(item, style='List Bullet')

    doc.add_heading('Segment Performance', level=1)
    doc.add_paragraph('The Cloud Platform segment continued to be the primary growth driver, generating $2.41B in revenue (50% of total) with 31.2% YoY growth. AI & Analytics contributed $1.21B (25%) growing 28.7%. Cybersecurity revenue was $723M (15%) with 15.4% growth. Professional Services generated $482M (10%) growing 8.1%.')

    doc.add_heading('Strategic Initiatives', level=1)
    doc.add_paragraph('Management outlined four strategic priorities for the remainder of FY2026:')
    for item in ['AI Platform Expansion: Accelerating enterprise AI adoption with new GenAI capabilities', 'Geographic Growth: Expanding sales teams in APAC (+40% headcount) and EMEA (+25%)', 'Cloud Migration: Converting 200+ on-premise customers to cloud platform', 'Cybersecurity M&A: Evaluating 3 acquisition targets to strengthen security portfolio', 'Partner Ecosystem: Launching 15 new ISV integrations by Q4']:
        doc.add_paragraph(item, style='List Bullet')

    doc.add_heading('Risk Assessment', level=1)
    doc.add_paragraph('The company identified several key risks: potential enterprise IT spending slowdown due to macroeconomic uncertainty, increasing competition from cloud hyperscalers, and approximately 150bps foreign currency headwind from USD strength. Additionally, customer concentration risk remains elevated with top 10 clients representing 28% of total revenue. Regulatory risk from the EU AI Act could impact product roadmap timelines.')

    doc.add_heading('Forward Guidance', level=1)
    doc.add_paragraph('Management raised full-year FY2026 guidance:')
    for item in ['Q2 2026 Revenue: $5.05B to $5.15B (21-23% YoY growth)', 'FY2026 Revenue: $20.5B to $21.0B (raised from prior $19.8B-$20.3B)', 'Operating Margin: Targeting 31-32% for Q2, expanding through year', 'EPS Guidance: $12.80 to $13.20 for full year', 'Capital Allocation: $2.5B share repurchase program authorized']:
        doc.add_paragraph(item, style='List Bullet')

    doc.add_heading('Analyst Consensus', level=1)
    doc.add_paragraph('Following the earnings release, 8 analysts raised price targets with an average target of $142 per share, representing 18% upside. The consensus rating moved to Strong Buy with 12 out of 14 analysts recommending the stock.')

    doc.save(path)
    print(f"Created: {path}")

create_sample_pptx(f"{VOLUME_PATH}/quarterly_review.pptx")
create_sample_docx(f"{VOLUME_PATH}/earnings_summary.docx")

# COMMAND ----------

# MAGIC %md
# MAGIC ## Step 2: Parse Documents with `ai_parse_document()`
# MAGIC
# MAGIC This is the key Databricks AI function — it reads raw file bytes and extracts structured text,
# MAGIC tables, and layout information. No external libraries needed.

# COMMAND ----------

# MAGIC %sql
# MAGIC -- Parse the PowerPoint presentation
# MAGIC SELECT
# MAGIC   'quarterly_review.pptx' AS file_name,
# MAGIC   ai_parse_document(
# MAGIC     content,
# MAGIC     map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING AS parsed_content
# MAGIC FROM read_files(
# MAGIC   '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/quarterly_review.pptx',
# MAGIC   format => 'binaryFile'
# MAGIC )

# COMMAND ----------

# MAGIC %sql
# MAGIC -- Parse the Word document
# MAGIC SELECT
# MAGIC   'earnings_summary.docx' AS file_name,
# MAGIC   ai_parse_document(
# MAGIC     content,
# MAGIC     map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING AS parsed_content
# MAGIC FROM read_files(
# MAGIC   '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/earnings_summary.docx',
# MAGIC   format => 'binaryFile'
# MAGIC )

# COMMAND ----------

# MAGIC %md
# MAGIC ## Step 3: Parse Both Documents and Compare with AI
# MAGIC
# MAGIC This single query parses both files and feeds them into `ai_query()` for comparison.

# COMMAND ----------

# MAGIC %sql
# MAGIC WITH pptx_parsed AS (
# MAGIC   SELECT ai_parse_document(
# MAGIC     content, map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING AS content
# MAGIC   FROM read_files(
# MAGIC     '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/quarterly_review.pptx',
# MAGIC     format => 'binaryFile'
# MAGIC   )
# MAGIC ),
# MAGIC doc_parsed AS (
# MAGIC   SELECT ai_parse_document(
# MAGIC     content, map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING AS content
# MAGIC   FROM read_files(
# MAGIC     '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/earnings_summary.docx',
# MAGIC     format => 'binaryFile'
# MAGIC   )
# MAGIC )
# MAGIC SELECT ai_query(
# MAGIC   'databricks-meta-llama-3-3-70b-instruct',
# MAGIC   CONCAT(
# MAGIC     'You are a document alignment analyst. Compare the following PowerPoint presentation and document.\n\n',
# MAGIC     '=== POWERPOINT PRESENTATION ===\n', pptx.content, '\n\n',
# MAGIC     '=== DOCUMENT ===\n', doc.content, '\n\n',
# MAGIC     'Produce a structured analysis:\n\n',
# MAGIC     '## Overall Similarity Score\n',
# MAGIC     'Rate alignment 0-100% with explanation.\n\n',
# MAGIC     '## Aligned Content\n',
# MAGIC     'Key topics consistent between both. Cite specific slides and sections.\n\n',
# MAGIC     '## Divergences Found\n',
# MAGIC     'Information that differs. For each state what PowerPoint says vs Document says, severity (Minor/Moderate/Major), and impact.\n\n',
# MAGIC     '## Content in PowerPoint Only\n',
# MAGIC     'Information in the presentation but missing from the document.\n\n',
# MAGIC     '## Content in Document Only\n',
# MAGIC     'Information in the document but missing from the presentation.\n\n',
# MAGIC     '## Recommendation\n',
# MAGIC     'What needs updating to bring documents into full alignment.\n\n',
# MAGIC     'Be specific with numbers, percentages, and exact claims.'
# MAGIC   )
# MAGIC ) AS comparison_report
# MAGIC FROM pptx_parsed pptx
# MAGIC CROSS JOIN doc_parsed doc

# COMMAND ----------

# MAGIC %md
# MAGIC ## Step 3b: Formatted Comparison Report
# MAGIC
# MAGIC Render the AI comparison as a nicely formatted, readable report:

# COMMAND ----------

# Fetch the latest comparison result and render it
report_df = spark.sql("""
WITH pptx_parsed AS (
  SELECT ai_parse_document(
    content, map('version', '2.0', 'descriptionElementTypes', '*')
  )::STRING AS content
  FROM read_files(
    '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/quarterly_review.pptx',
    format => 'binaryFile'
  )
),
doc_parsed AS (
  SELECT ai_parse_document(
    content, map('version', '2.0', 'descriptionElementTypes', '*')
  )::STRING AS content
  FROM read_files(
    '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/earnings_summary.docx',
    format => 'binaryFile'
  )
)
SELECT
  pptx.content AS pptx_text,
  doc.content AS doc_text,
  ai_query(
    'databricks-meta-llama-3-3-70b-instruct',
    CONCAT(
      'You are a document alignment analyst. Compare the following PowerPoint presentation and document.\\n\\n',
      '=== POWERPOINT PRESENTATION ===\\n', pptx.content, '\\n\\n',
      '=== DOCUMENT ===\\n', doc.content, '\\n\\n',
      'Produce a structured analysis:\\n\\n',
      '## Overall Similarity Score\\n',
      'Rate alignment 0-100% with explanation.\\n\\n',
      '## Aligned Content\\n',
      'Key topics consistent between both. Cite specific slides and sections.\\n\\n',
      '## Divergences Found\\n',
      'Information that differs. For each state what PowerPoint says vs Document says, severity (Minor/Moderate/Major), and impact.\\n\\n',
      '## Content in PowerPoint Only\\n',
      'Information in the presentation but missing from the document.\\n\\n',
      '## Content in Document Only\\n',
      'Information in the document but missing from the presentation.\\n\\n',
      '## Recommendation\\n',
      'What needs updating to bring documents into full alignment.\\n\\n',
      'Be specific with numbers, percentages, and exact claims.'
    )
  ) AS comparison_report
FROM pptx_parsed pptx
CROSS JOIN doc_parsed doc
""")

row = report_df.collect()[0]
report_text = row["comparison_report"]
pptx_len = len(row["pptx_text"])
doc_len = len(row["doc_text"])

# Build formatted HTML report
html = f"""
<div style="font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 900px; margin: 0 auto; color: #e2e8f0;">

  <div style="background: linear-gradient(135deg, #1e3a5f 0%, #0a1628 100%); border-radius: 12px; padding: 24px 32px; margin-bottom: 24px;">
    <h1 style="margin: 0 0 8px 0; font-size: 22px; color: #f8fafc;">Document Alignment Report</h1>
    <p style="margin: 0; font-size: 13px; color: #94a3b8;">AI-powered comparison using Databricks <code style="background: rgba(255,255,255,0.1); padding: 2px 6px; border-radius: 4px; font-size: 12px;">ai_parse_document()</code> + <code style="background: rgba(255,255,255,0.1); padding: 2px 6px; border-radius: 4px; font-size: 12px;">ai_query()</code></p>
  </div>

  <div style="display: flex; gap: 16px; margin-bottom: 24px;">
    <div style="flex: 1; background: #1e293b; border: 1px solid #334155; border-radius: 10px; padding: 16px;">
      <div style="font-size: 11px; color: #64748b; text-transform: uppercase; letter-spacing: 0.5px;">PowerPoint</div>
      <div style="font-size: 15px; font-weight: 600; color: #fb923c; margin-top: 4px;">quarterly_review.pptx</div>
      <div style="font-size: 12px; color: #94a3b8; margin-top: 4px;">{pptx_len:,} characters parsed</div>
    </div>
    <div style="flex: 0; display: flex; align-items: center; font-size: 20px; color: #3b82f6;">⇄</div>
    <div style="flex: 1; background: #1e293b; border: 1px solid #334155; border-radius: 10px; padding: 16px;">
      <div style="font-size: 11px; color: #64748b; text-transform: uppercase; letter-spacing: 0.5px;">Document</div>
      <div style="font-size: 15px; font-weight: 600; color: #60a5fa; margin-top: 4px;">earnings_summary.docx</div>
      <div style="font-size: 12px; color: #94a3b8; margin-top: 4px;">{doc_len:,} characters parsed</div>
    </div>
  </div>

  <div style="background: #1e293b; border: 1px solid #334155; border-radius: 10px; padding: 24px 32px;">
"""

# Convert markdown to styled HTML
import re

lines = report_text.split("\n")
for line in lines:
    stripped = line.strip()
    if not stripped:
        html += "<br/>"
    elif stripped.startswith("## "):
        heading = stripped[3:]
        # Color-code specific sections
        if "similarity" in heading.lower() or "score" in heading.lower():
            color = "#22c55e"
            icon = "🎯"
        elif "aligned" in heading.lower():
            color = "#3b82f6"
            icon = "✅"
        elif "divergen" in heading.lower():
            color = "#f59e0b"
            icon = "⚠️"
        elif "powerpoint only" in heading.lower():
            color = "#fb923c"
            icon = "📊"
        elif "document only" in heading.lower():
            color = "#60a5fa"
            icon = "📄"
        elif "recommend" in heading.lower():
            color = "#a78bfa"
            icon = "💡"
        else:
            color = "#e2e8f0"
            icon = "📋"
        html += f'<h2 style="font-size: 16px; color: {color}; margin: 20px 0 10px 0; padding-bottom: 8px; border-bottom: 1px solid #334155;">{icon} {heading}</h2>'
    elif stripped.startswith("### "):
        html += f'<h3 style="font-size: 14px; color: #cbd5e1; margin: 14px 0 6px 0;">{stripped[4:]}</h3>'
    elif stripped.startswith("- **") or stripped.startswith("* **"):
        # Bold list item
        content = stripped[2:]
        content = re.sub(r'\*\*(.*?)\*\*', r'<strong style="color: #f8fafc;">\1</strong>', content)
        html += f'<div style="font-size: 13px; color: #94a3b8; padding: 4px 0 4px 16px; border-left: 2px solid #334155; margin: 4px 0;">• {content}</div>'
    elif stripped.startswith("- ") or stripped.startswith("* "):
        content = stripped[2:]
        content = re.sub(r'\*\*(.*?)\*\*', r'<strong style="color: #f8fafc;">\1</strong>', content)
        html += f'<div style="font-size: 13px; color: #94a3b8; padding: 4px 0 4px 16px; border-left: 2px solid #334155; margin: 4px 0;">• {content}</div>'
    elif re.match(r'^\d+\.', stripped):
        content = re.sub(r'^\d+\.\s*', '', stripped)
        content = re.sub(r'\*\*(.*?)\*\*', r'<strong style="color: #f8fafc;">\1</strong>', content)
        html += f'<div style="font-size: 13px; color: #94a3b8; padding: 4px 0 4px 16px; margin: 4px 0;">{stripped[:2]} {content[len(stripped[:2]):]}</div>'
    else:
        content = re.sub(r'\*\*(.*?)\*\*', r'<strong style="color: #f8fafc;">\1</strong>', stripped)
        html += f'<p style="font-size: 13px; color: #cbd5e1; margin: 6px 0; line-height: 1.6;">{content}</p>'

html += """
  </div>

  <div style="text-align: center; margin-top: 16px; font-size: 11px; color: #475569;">
    Generated by Databricks AI Functions — ai_parse_document() + ai_query()
  </div>
</div>
"""

displayHTML(html)

# COMMAND ----------

# MAGIC %md
# MAGIC ## Step 4: Save Report to Delta Table

# COMMAND ----------

# MAGIC %sql
# MAGIC CREATE TABLE IF NOT EXISTS lakedoom_demo_catalog.capital_markets_ai.document_comparisons (
# MAGIC   comparison_id STRING,
# MAGIC   pptx_file STRING,
# MAGIC   doc_file STRING,
# MAGIC   pptx_parsed_text STRING,
# MAGIC   doc_parsed_text STRING,
# MAGIC   comparison_report STRING,
# MAGIC   created_at TIMESTAMP
# MAGIC );
# MAGIC
# MAGIC WITH pptx_parsed AS (
# MAGIC   SELECT ai_parse_document(
# MAGIC     content, map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING AS content
# MAGIC   FROM read_files(
# MAGIC     '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/quarterly_review.pptx',
# MAGIC     format => 'binaryFile'
# MAGIC   )
# MAGIC ),
# MAGIC doc_parsed AS (
# MAGIC   SELECT ai_parse_document(
# MAGIC     content, map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING AS content
# MAGIC   FROM read_files(
# MAGIC     '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/earnings_summary.docx',
# MAGIC     format => 'binaryFile'
# MAGIC   )
# MAGIC )
# MAGIC INSERT INTO lakedoom_demo_catalog.capital_markets_ai.document_comparisons
# MAGIC SELECT
# MAGIC   CONCAT('CMP-', date_format(current_timestamp(), 'yyyyMMddHHmmss')) AS comparison_id,
# MAGIC   'quarterly_review.pptx' AS pptx_file,
# MAGIC   'earnings_summary.docx' AS doc_file,
# MAGIC   pptx.content AS pptx_parsed_text,
# MAGIC   doc.content AS doc_parsed_text,
# MAGIC   ai_query(
# MAGIC     'databricks-meta-llama-3-3-70b-instruct',
# MAGIC     CONCAT(
# MAGIC       'Compare the following PowerPoint and document. Provide: Overall Similarity Score (0-100%), Aligned Content, Divergences (with severity), Content in PowerPoint Only, Content in Document Only, and Recommendations. Be specific with numbers.\n\n',
# MAGIC       '=== POWERPOINT ===\n', pptx.content, '\n\n=== DOCUMENT ===\n', doc.content
# MAGIC     )
# MAGIC   ) AS comparison_report,
# MAGIC   current_timestamp() AS created_at
# MAGIC FROM pptx_parsed pptx
# MAGIC CROSS JOIN doc_parsed doc

# COMMAND ----------

# MAGIC %sql
# MAGIC -- View saved comparison reports
# MAGIC SELECT comparison_id, pptx_file, doc_file,
# MAGIC        LENGTH(pptx_parsed_text) AS pptx_chars,
# MAGIC        LENGTH(doc_parsed_text) AS doc_chars,
# MAGIC        LEFT(comparison_report, 500) AS report_preview,
# MAGIC        created_at
# MAGIC FROM lakedoom_demo_catalog.capital_markets_ai.document_comparisons
# MAGIC ORDER BY created_at DESC

# COMMAND ----------

# MAGIC %md
# MAGIC ## Step 5: Compare ANY Documents at Scale
# MAGIC
# MAGIC Drop multiple file pairs into the volume and compare them all in one query:

# COMMAND ----------

# MAGIC %sql
# MAGIC -- Parse ALL documents in the volume at once
# MAGIC SELECT
# MAGIC   regexp_extract(path, '([^/]+)$', 1) AS filename,
# MAGIC   LENGTH(ai_parse_document(
# MAGIC     content, map('version', '2.0', 'descriptionElementTypes', '*')
# MAGIC   )::STRING) AS parsed_length_chars
# MAGIC FROM read_files(
# MAGIC   '/Volumes/lakedoom_demo_catalog/capital_markets_ai/document_comparison/',
# MAGIC   format => 'binaryFile'
# MAGIC )

# COMMAND ----------

# MAGIC %md
# MAGIC ## Summary
# MAGIC
# MAGIC | Step | Function | What It Does |
# MAGIC |------|----------|-------------|
# MAGIC | **Parse** | `ai_parse_document()` | Converts raw .pptx/.docx bytes into structured text — handles slides, tables, headers, bullets |
# MAGIC | **Compare** | `ai_query()` | Uses a foundation model to analyze alignment between two parsed documents |
# MAGIC | **Persist** | `INSERT INTO` Delta | Saves parsed text and comparison reports to a governed table |
# MAGIC | **Scale** | `read_files()` + batch | Process entire document libraries in a single SQL query |
# MAGIC
# MAGIC **Key advantage:** Zero Python library installs. Everything runs as SQL functions on the SQL Warehouse.
# MAGIC The only library install (Step 1) is for generating sample demo files — not needed for the actual comparison pipeline.
